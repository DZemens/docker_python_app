import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

class PPT_Creator(object):
    def __init__(self):
        self._data_frames = []
    @property
    def data_frames(self):
        return self._data_frames
    def add_frame(self, df):
        self._data_frames.append(df)
    def create(self):
        data_frames = self.data_frames
        p = Presentation()
        layout = p.slide_layouts[6]
        for frame in data_frames:
            # slide/chart with six series
            cd = CategoryChartData()
            cd.categories = frame.index
            for name in frame.columns:
                cd.add_series(name, frame[name], '0%')
            slide = p.slides.add_slide( layout )
            shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, 0, 0, 9143301, 6158000, cd)
            # slide/chart with five series
            cd = CategoryChartData()
            cd.categories = frame.index
            for name in frame.columns[:-1]:
                cd.add_series(name, frame[name], '0%')
            slide = p.slides.add_slide( layout )
            shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, 0, 0, 9143301, 6158000, cd)
        return p